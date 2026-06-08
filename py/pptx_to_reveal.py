#!/usr/bin/env python3
"""Convert .pptx files to self-contained reveal.js HTML presentations."""

import argparse
import base64
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

TITLE_PLACEHOLDER_TYPES = {0, 3}  # idx 0 = title, 3 = center_title (PP_PLACEHOLDER values)

REVEAL_CDN = "https://cdn.jsdelivr.net/npm/reveal.js@5"

HTML_TEMPLATE = """\
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<link rel="stylesheet" href="{cdn}/dist/reveal.css">
<link rel="stylesheet" href="{cdn}/dist/theme/black.css">
<style>
  .reveal .slides section {{ text-align: left; }}
  .reveal .slides section h2 {{ text-align: center; }}
  .reveal ul {{ margin-left: 1em; }}
  .reveal img {{ max-height: 60vh; max-width: 90%; display: block; margin: 0 auto; }}
  .reveal pre {{ font-size: 0.55em; }}
  .slide-number-in {{ color: #aaa; font-size: 0.4em; }}
</style>
</head>
<body>
<div class="reveal">
  <div class="slides">
{slides}
  </div>
</div>
<script src="{cdn}/dist/reveal.js"></script>
<script>
  Reveal.initialize({{
    hash: true,
    slideNumber: 'c/t',
    transition: 'slide',
    plugins: []
  }});
</script>
</body>
</html>
"""

SLIDE_TEMPLATE = """\
    <section>
{content}
    </section>"""


def _pt_to_em(pt_value):
    if pt_value is None:
        return None
    try:
        return round(Pt(1).pt and pt_value.pt / 16, 2)
    except Exception:
        return None


def _placeholder_format(shape):
    """Return placeholder_format if shape is a placeholder, else None."""
    if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
        return None
    try:
        return shape.placeholder_format
    except Exception:
        return None


def _is_title(shape) -> bool:
    pf = _placeholder_format(shape)
    if pf is None:
        return False
    return pf.idx == 0 or int(pf.type) in TITLE_PLACEHOLDER_TYPES


def _shape_to_html(shape) -> str:
    parts = []

    # Images
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img_bytes = shape.image.blob
            mime = shape.image.content_type or "image/png"
            b64 = base64.b64encode(img_bytes).decode()
            parts.append(f'      <img src="data:{mime};base64,{b64}" alt="">')
        except Exception:
            pass
        return "\n".join(parts)

    # Text frames
    if not shape.has_text_frame:
        return ""

    tf = shape.text_frame
    is_title = _is_title(shape)

    paragraphs_html = []
    for para in tf.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Inline runs with formatting
        inline = ""
        for run in para.runs:
            run_text = run.text
            if not run_text:
                continue
            run_text = run_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            style_parts = []
            try:
                if run.font.bold:
                    run_text = f"<strong>{run_text}</strong>"
                if run.font.italic:
                    run_text = f"<em>{run_text}</em>"
                if run.font.underline:
                    run_text = f"<u>{run_text}</u>"
                size_em = _pt_to_em(run.font.size)
                if size_em:
                    style_parts.append(f"font-size:{size_em}em")
            except Exception:
                pass
            if style_parts:
                run_text = f'<span style="{"; ".join(style_parts)}">{run_text}</span>'
            inline += run_text

        if not inline:
            inline = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        level = para.level or 0
        indent = "  " * level

        if is_title:
            paragraphs_html.append(f"      <h2>{inline}</h2>")
        else:
            paragraphs_html.append(f"      {indent}<li>{inline}</li>")

    if not paragraphs_html:
        return ""

    if is_title:
        return "\n".join(paragraphs_html)

    # Wrap non-title content in a <ul>
    return "      <ul>\n" + "\n".join(paragraphs_html) + "\n      </ul>"


def slide_to_html(slide, index: int) -> str:
    content_parts = []

    # Sort shapes: title-like placeholders first, then others by position
    def shape_order(s):
        pf = _placeholder_format(s)
        if pf is not None:
            return (0, pf.idx)
        return (1, getattr(s, "top", 0) or 0)

    for shape in sorted(slide.shapes, key=shape_order):
        html = _shape_to_html(shape)
        if html:
            content_parts.append(html)

    # Speaker notes
    try:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            notes_escaped = notes_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            content_parts.append(f'      <aside class="notes">{notes_escaped}</aside>')
    except Exception:
        pass

    if not content_parts:
        content_parts.append("      <!-- empty slide -->")

    return SLIDE_TEMPLATE.format(content="\n".join(content_parts))


def convert(pptx_path: Path, output_path: Path | None = None) -> Path:
    prs = Presentation(pptx_path)

    slides_html = []
    for i, slide in enumerate(prs.slides):
        slides_html.append(slide_to_html(slide, i + 1))

    # Derive title from first non-empty title shape
    title = pptx_path.stem
    for slide in prs.slides:
        for shape in slide.shapes:
            if _is_title(shape) and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    title = t
                    break
        else:
            continue
        break

    html = HTML_TEMPLATE.format(
        title=title,
        cdn=REVEAL_CDN,
        slides="\n".join(slides_html),
    )

    out = output_path or pptx_path.with_suffix(".html")
    out.write_text(html, encoding="utf-8")
    return out


def main():
    parser = argparse.ArgumentParser(description="Convert .pptx to reveal.js HTML")
    parser.add_argument("files", nargs="*", help=".pptx files (omit to convert all in docs/)")
    parser.add_argument("--root", default=".", help="Repo root (default: current dir)")
    args = parser.parse_args()

    root = Path(args.root)

    EXCLUDE_DIRS = {".venv", "site", ".git"}

    if args.files:
        paths = [Path(f) for f in args.files]
    else:
        paths = [
            p for p in root.rglob("*.pptx")
            if not any(part in EXCLUDE_DIRS for part in p.parts)
        ]

    if not paths:
        print("No .pptx files found.")
        sys.exit(0)

    for p in paths:
        try:
            out = convert(p)
            print(f"  OK  {p.relative_to(root)} -> {out.name}")
        except Exception as e:
            print(f"  ERR {p}: {e}")


if __name__ == "__main__":
    main()
